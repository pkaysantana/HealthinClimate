"""Build the HeatGuard pitch deck on top of the London Hackathon 2026 template.

Reuses the template's theme (Arial, the accent palette) and its named layouts
(TITLE / TITLE_AND_BODY / BIG_NUMBER), fills in HeatGuard content following the
required structure (Problem / Solution / How it works / Why meaningful / Future
support / Demo), and embeds dashboard screenshots so the .pptx is self-contained.

Usage:
    python scripts/build_deck.py [TEMPLATE.pptx] [OUTPUT.pptx]
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(sys.argv[1]) if len(sys.argv) > 1 else next(
    ROOT.glob(".context/attachments/**/*.pptx"), ROOT / "template.pptx"
)
OUT = Path(sys.argv[2]) if len(sys.argv) > 2 else ROOT / "docs" / "HeatGuard_London_Hackathon_2026.pptx"
SHOTS = ROOT / ".context" / "screenshots"

ORANGE = RGBColor(0xFF, 0x6A, 0x00)   # heat accent
RED = RGBColor(0xDC, 0x26, 0x26)
TEAL = RGBColor(0x00, 0x97, 0xA7)
DARK = RGBColor(0x21, 0x21, 0x21)
GREY = RGBColor(0x59, 0x59, 0x59)

prs = Presentation(str(TEMPLATE))
LAYOUTS = {l.name: l for l in prs.slide_masters[0].slide_layouts}

# remove the template's example slides (drop the relationship too, so the orphaned
# slide parts are not serialized — otherwise duplicate slideN.xml names appear)
_sldIdLst = prs.slides._sldIdLst
for sldId in list(_sldIdLst):
    rId = sldId.get(qn("r:id"))
    _sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def add(layout: str):
    return prs.slides.add_slide(LAYOUTS[layout])


def _set(ph, lines, *, size=None, bold=None, color=None):
    """Fill a placeholder. `lines` is a str, a list of str, or (text, level) tuples."""
    if isinstance(lines, str):
        lines = [lines]
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        text, level = (line if isinstance(line, tuple) else (line, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = "Arial"
        if size:
            f.size = Pt(size - 2 * level)
        if bold is not None:
            f.bold = bold
        if color is not None:
            f.color.rgb = color


def title_body(title, body, *, body_size=13):
    s = add("TITLE_AND_BODY")
    _set(s.placeholders[0], title, size=26, bold=True, color=DARK)
    _set(s.placeholders[1], body, size=body_size)
    s.placeholders[1].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return s


def big_number(number, caption, *, num_color=RED):
    s = add("BIG_NUMBER")
    _set(s.placeholders[0], number, size=96, bold=True, color=num_color)
    _set(s.placeholders[1], caption, size=15, color=GREY)
    s.placeholders[1].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return s


def picture(slide, name, left, top, width):
    f = SHOTS / name
    if f.exists():
        slide.shapes.add_picture(str(f), Inches(left), Inches(top), width=Inches(width))


# ---- 1. Title --------------------------------------------------------------
s = add("TITLE")
_set(s.placeholders[0], "HeatGuard", size=54, bold=True, color=ORANGE)
_set(s.placeholders[1], [
    "Adaptive heat-safety scheduling for Gulf outdoor crews",
    "Replacing the blunt calendar ban with a condition-responsive, provable system",
    "London Hackathon 2026",
], size=16, color=DARK)

# ---- 2. Team ---------------------------------------------------------------
title_body("Team", [
    "[Add each member's name and one-line skillset]",
    ("A multidisciplinary build: thermal-physiology + occupational-health standards, "
     "backend/data engineering, and front-end / product.", 0),
])

# ---- 3. Problem ------------------------------------------------------------
title_body("What is the problem?", [
    "Millions of migrant workers do outdoor labour across the Gulf, where summer heat exceeds the limits of human thermoregulation.",
    ("24,000,000 migrant workers in the Arab States — 41.4%, the highest share of any world region.", 0),
    ("~10,000 migrant deaths/year — and over half are certified with no underlying cause, so the heat toll is hidden.", 0),
    ("The control is a fixed CALENDAR ban (e.g. Saudi 12:00–15:00, 15 Jun–15 Sep). It is wrong in BOTH directions:", 0),
    ("misses May/September heat, humid mornings, and unacclimatized newcomers — the people who actually die", 1),
    ("and needlessly stops safe work, so employers resent and evade it.", 1),
], body_size=13)

# ---- 4. The danger (big number) -------------------------------------------
big_number("~10,000", [
    "migrant deaths a year in the Gulf (all-cause). Most are certified without an underlying cause, "
    "so the heat-attributable toll is hidden — and only 1 of 19 worker-heat studies worldwide comes "
    "from the Gulf. The danger is huge and largely unmeasured.",
])

# ---- 5. Solution -----------------------------------------------------------
s = title_body("What is your solution?", [
    "A site-level, WBGT-driven work–rest–hydration scheduler for the supervisor — one cheap sensor per site, not a wearable per worker.",
    ("Senses WBGT (or a ~$300 on-site meter) → outputs the actual ACGIH/ISO 7243 work-rest cycle and ISO 7933 hydration target for current conditions and the worker.", 0),
    ("Broadcasts one signal everyone understands: WORK · REST IN SHADE · DRINK NOW · STOP.", 0),
    ("Ramps new arrivals in safely (NIOSH) and logs everything to a tamper-evident, privacy-by-design worker-protection record.", 0),
], body_size=12)
# narrow the body and add the timeline screenshot on the right
s.placeholders[1].width = Inches(5.0)
picture(s, "03_timeline.png", left=5.2, top=1.5, width=4.5)

# ---- 6. How it works -------------------------------------------------------
title_body("How does it work?", [
    "Data: free Open-Meteo reanalysis weather (cached for an offline demo).",
    ("WBGT: the validated Liljegren model (ECMWF thermofeel) with a vendored solar calc — or a measured on-site reading.", 0),
    ("Schedule: hardcoded ACGIH TLV / Action-Limit tables (ISO 7243); ISO 7933 Predicted Heat Strain (pythermalcomfort) for hydration and max exposure; NIOSH acclimatization ramp.", 0),
    ("Decision: the most-conservative of table · ramp · physiology → one signal. Every reading is SHA-256 hash-chained.", 0),
    ("One pure Python engine with a FastAPI API, a React dashboard, a Streamlit app, and a validation notebook — 79 tests + green CI.", 0),
], body_size=12)

# ---- 7. Why meaningful -----------------------------------------------------
title_body("Why is it meaningful?", [
    "It makes a cheap, proven intervention adaptive and auditable — the missing implementation layer.",
    ("Per 100-worker crew, one season: 1,237 dangerous hours the ban missed; 7.7 acute-kidney-injury cases averted; productivity maintained or raised 10–20%.", 0),
    ("Validated: the impact model reproduces the real La Isla / Adelante (Nicaragua) outcome — AKI −94%, productivity +10–20%.", 0),
    ("Productivity-positive with a ~6-week payback, plus a compliance shield — so it actually gets switched on.", 0),
], body_size=13)

# ---- 8. Lives saved (big number) ------------------------------------------
big_number("~1,900", [
    "lives saved across a 5,000,000-worker regional season — with ~383,000 kidney-injury cases averted "
    "and $1.6–2.5B in value (illustrative, conservative). A 100,000-worker megaproject: ~38 lives, "
    "7,670 AKI cases averted, $32–50M of value.",
], num_color=ORANGE)

# ---- 9. The business case (big number) ------------------------------------
big_number("3–5×", [
    "ROI with a ~6-week payback in Dubai; 7–10× in Riyadh, where the blunt ban also destroys safe work. "
    "About $95/worker, mostly one-time capital. Not a safety cost — a productivity gain.",
], num_color=TEAL)

# ---- 10. Future support ----------------------------------------------------
title_body("What would you do with more support?", [
    "Pilot with a Gulf labour-supply contractor: deploy cheap on-site WBGT meters + a site horn/light and validate against real outcomes.",
    ("Development-finance pathway: lenders and ESG mandates already attach occupational-safety conditions — HeatGuard is the affordable implementation-and-verification layer.", 0),
    ("Close the data gap: there is almost no public Gulf worker-heat data — partner to collect it.", 0),
    ("Scale the engine to other heat-exposed sectors and regions (agriculture, logistics).", 0),
], body_size=13)

# ---- 11. Product demo ------------------------------------------------------
s = title_body("Product demo", [
    "Live dashboard: the signal, the calendar-ban-vs-HeatGuard timeline, sensor-vs-estimate, per-worker individualization, the worker-protection record, and the business case.",
    ("github.com/fcistud/heatguard   ·   79 tests + green CI   ·   MIT", 0),
    ("A prototype, not certified safety equipment.", 0),
], body_size=12)
s.placeholders[1].width = Inches(4.7)
picture(s, "04_impact.png", left=5.0, top=1.4, width=4.7)
picture(s, "10_measured.png", left=5.0, top=3.3, width=4.7)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
